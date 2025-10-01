import os
import json
import time
import uuid
import glob
import fitz  # PDF
import requests
from datetime import datetime

from fastapi import FastAPI, UploadFile, Form, File, Header, HTTPException
from fastapi.responses import JSONResponse, StreamingResponse, PlainTextResponse
from fastapi.middleware.cors import CORSMiddleware
from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
from pydantic import BaseModel, Field
from typing import Optional, List, Dict, Any

from openai import OpenAI
from dotenv import load_dotenv
from database import get_data_by_request_id, update_faq_result, insert_full_record

import jwt  # pyjwt

# ──────────────────────────────────────────────────────────────────────────────
# Setup
# ──────────────────────────────────────────────────────────────────────────────
load_dotenv()
client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
JWT_SECRET = os.getenv("JWT_SECRET", "").strip()  # <-- set this in env
if not JWT_SECRET:
    # Safe fallback for dev only (do NOT use in prod)
    JWT_SECRET = "dev-only-secret-change-me"

app = FastAPI()
UPLOAD_DIR = "./uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # set to your domain(s) in prod
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ──────────────────────────────────────────────────────────────────────────────
# Existing helpers (unchanged)
# ──────────────────────────────────────────────────────────────────────────────
def extract_text_from_pdf(pdf_path):
    text = ""
    with fitz.open(pdf_path) as doc:
        for page in doc:
            text += page.get_text()
    return text

def extract_text_from_docx(docx_path):
    doc = Document(docx_path)
    return "\n".join(para.text.strip() for para in doc.paragraphs if para.text.strip())

def extract_text_from_pptx(pptx_path):
    text = ""
    prs = Presentation(pptx_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_url(url):
    try:
        response = requests.get(url, timeout=30)
        soup = BeautifulSoup(response.content, "html.parser")
        for tag in soup(["script", "style"]):
            tag.decompose()
        return ' '.join(soup.stripped_strings)
    except Exception:
        return ""

def generate_questions_and_answers(text, question_number, questions, faq_examples):
    examples_text = ""
    for idx, item in enumerate(faq_examples[:5], 1):
        examples_text += f"{idx}. س: {item['question']}\n   ج: {item['answer']}\n"

    prompt = f"""
نمط الأسئلة الشائعة لدينا كالتالي:
{examples_text}

اقرأ النص التالي واستخرج {question_number} سؤالًا شائعًا مع إجابته بطريقة احترافية:

\"\"\"{text}\"\"\"

مع الإجابة عن هذه الأسئلة:
{questions}
"""
    completion = client.chat.completions.create(
        model="gpt-5-mini",
        messages=[{"role": "user", "content": prompt}]
    )
    return completion.choices[0].message.content

# ──────────────────────────────────────────────────────────────────────────────
# Existing endpoint (unchanged)
# ──────────────────────────────────────────────────────────────────────────────
@app.post("/process-faq")
async def process_faq(
    file: UploadFile = File(None),
    url: str = Form(None),
    questions_number: int = Form(10),
    custom_questions: str = Form(""),
    user_id: int = Form(...)
):
    try:
        file_path = ""
        if file:
            filename = f"{int(datetime.now().timestamp())}_{file.filename}"
            file_path = os.path.join(UPLOAD_DIR, filename)
            with open(file_path, "wb") as f:
                f.write(await file.read())
        elif url:
            file_path = None
        else:
            return JSONResponse({"error": "رابط أو ملف مطلوب"}, status_code=400)

        # Extract
        if file_path:
            ext = os.path.splitext(file_path)[1].lower()
            if ext == ".pdf":
                extracted_text = extract_text_from_pdf(file_path)
            elif ext in [".doc", ".docx"]:
                extracted_text = extract_text_from_docx(file_path)
            elif ext == ".pptx":
                extracted_text = extract_text_from_pptx(file_path)
            else:
                return JSONResponse({"error": "نوع الملف غير مدعوم"}, status_code=400)
        else:
            extracted_text = extract_text_from_url(url)

        if not extracted_text.strip():
            return JSONResponse({"error": "لا يوجد محتوى صالح للمعالجة"}, status_code=400)

        with open("faq_examples.json", "r", encoding="utf-8") as f:
            faq_examples = json.load(f)

        # LLM
        faq_result = generate_questions_and_answers(
            extracted_text, questions_number, custom_questions, faq_examples
        )

        # Save DB
        saved = insert_full_record(user_id, file_path, url, questions_number, custom_questions, faq_result)

        # Cleanup
        if file_path and os.path.exists(file_path):
            os.remove(file_path)
            cleanup_old_files(days=7)

        return JSONResponse(content={"questions_and_answers": faq_result})

    except Exception as e:
        import traceback
        traceback.print_exc()
        return JSONResponse({"error": str(e)}, status_code=500)

def cleanup_old_files(days=7):
    now = time.time()
    for f in glob.glob(f"{UPLOAD_DIR}/*"):
        if os.path.isfile(f) and os.stat(f).st_mtime < now - days * 86400:
            os.remove(f)

# ──────────────────────────────────────────────────────────────────────────────
# NEW: Chat models + helpers
# ──────────────────────────────────────────────────────────────────────────────
class SessionIn(BaseModel):
    user_id: int
    wp_nonce: Optional[str] = None

class SessionOut(BaseModel):
    session_id: str
    token: str

class VisibleValue(BaseModel):
    # We accept anything the WP plugin sends (all optional)
    id: Optional[int] = None
    source_url: Optional[str] = None
    file_name: Optional[str] = None
    questions_number: Optional[int] = None
    custom_questions: Optional[str] = None
    last_result: Optional[str] = None
    created_at: Optional[str] = None
    updated_at: Optional[str] = None
    user_id: Optional[int] = None
    # from localStorage
    article: Optional[str] = None

class ChatIn(BaseModel):
    session_id: str
    user_id: int
    message: str
    visible_values: List[VisibleValue] = Field(default_factory=list)

def _make_jwt(session_id: str, user_id: int) -> str:
    payload = {
        "sid": session_id,
        "uid": user_id,
        "iat": int(time.time()),
        "exp": int(time.time()) + 60 * 60 * 2,  # 2 hours
    }
    return jwt.encode(payload, JWT_SECRET, algorithm="HS256")

def _verify_jwt(bearer: Optional[str]):
    if not bearer or not bearer.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Missing Authorization header")
    token = bearer.split(" ", 1)[1]
    try:
        jwt.decode(token, JWT_SECRET, algorithms=["HS256"])
    except jwt.InvalidTokenError:
        raise HTTPException(status_code=401, detail="Invalid token")

def _values_to_context(values: List[VisibleValue]) -> str:
    """Builds a concise Arabic context block from visible values."""
    if not values:
        return "لا توجد بيانات مرئية حالياً لهذا المستخدم."
    v = values[0]
    parts = []
    if v.source_url:        parts.append(f"المصدر: {v.source_url}")
    elif v.file_name:       parts.append(f"الملف: {v.file_name}")
    if v.questions_number is not None: parts.append(f"عدد الأسئلة: {v.questions_number}")
    if v.custom_questions:  parts.append(f"أسئلة مخصصة: {v.custom_questions[:200]}")
    # prefer article (fresh), else last_result
    article_text = v.article or v.last_result
    if article_text:
        parts.append(f"أحدث ناتج/مقال: {article_text[:800]}")
    if v.created_at: parts.append(f"تاريخ الإنشاء: {v.created_at}")
    if v.updated_at: parts.append(f"آخر تحديث: {v.updated_at}")
    return " | ".join(parts) if parts else "لا توجد تفاصيل كافية."

# ──────────────────────────────────────────────────────────────────────────────
# NEW: Chat routes
# ──────────────────────────────────────────────────────────────────────────────
@app.post("/session", response_model=SessionOut)
def create_session(body: SessionIn):
    sid = str(uuid.uuid4())
    token = _make_jwt(sid, body.user_id)
    return SessionOut(session_id=sid, token=token)

@app.post("/chat")
def chat(body: ChatIn, authorization: Optional[str] = Header(None)):
    _verify_jwt(authorization)

    context = _values_to_context(body.visible_values)
    sys_prompt = (
        "أنت مساعد يجيب باحتراف اعتمادًا على بيانات المستخدم المرئية أدناه. "
        "إن لم تتوفر المعلومة في البيانات المتاحة، قل ذلك صراحةً واقترح طريقة للحصول عليها.\n\n"
        f"البيانات المرئية الحالية:\n{context}\n\n"
        "عند تقديم الإجابة:\n"
        "- كن موجزًا ومباشرًا مع نقاط مرتبة عند الحاجة.\n"
        "- لا تفترض معلومات غير موجودة في البيانات أعلاه.\n"
        "- إذا سُئلت عن محتوى الأسئلة/الأجوبة الأخيرة، فاستند إلى أحدث ناتج/مقال إن وُجد."
    )

    user_msg = body.message or ""

    def stream():
        try:
            response = client.chat.completions.create(
                model="gpt-5-mini",
                messages=[
                    {"role": "system", "content": sys_prompt},
                    {"role": "user",   "content": user_msg}
                ],
                stream=True
            )
            for chunk in response:
                if chunk.choices and getattr(chunk.choices[0].delta, "content", None):
                    yield chunk.choices[0].delta.content
        except Exception as e:
            # Fallback: non-stream error text
            yield f"\n[خطأ في البث: {type(e).__name__}: {e}]"

    # Try streaming; if client/infra blocks streaming, caller will still get text/plain
    return StreamingResponse(stream(), media_type="text/plain")



